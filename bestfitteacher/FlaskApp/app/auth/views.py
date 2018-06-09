from datetime import datetime
import re
from flask import render_template, redirect, request, url_for,session, flash,jsonify,make_response
from PyPDF2 import PdfFileReader, PdfFileWriter
from pathlib import Path
from functools import wraps, update_wrapper
import csv
from flask_login import login_user, logout_user, login_required, \
    current_user
from . import auth
import  pandas as pd
from functools import wraps
from .. import db
from ..models import User,Book
from ..email import send_email
from .forms import LoginForm, RegistrationForm, ChangePasswordForm,\
    PasswordResetRequestForm, PasswordResetForm, ChangeEmailForm
import uuid
import gc
from pptx import Presentation
import pdfkit
import json
from django.utils.encoding import smart_str, smart_unicode
from bs4 import BeautifulSoup
import urllib2
from jinja2 import Environment, FileSystemLoader
import os
import subprocess
from django.core.validators import URLValidator
import validators
import glob
class Flashcardqa:
    questionnum= 0
    question ="question"
    answer= "answer"
from django.core.exceptions import ValidationError
# THIS_DIR = os.path.dirname(os.path.abspath(__file__))
# THIS_DIR = '/home/atul/Documents/SOFTWARE/OwnPrograms/APPLICATIONS/tbook2wbook/atul_flasky/app'
THIS_DIR='/var/www/FlaskApp/FlaskApp' +'/app'
stringdirectory=THIS_DIR + '/templates/toinsert.html'
stringtoinsert = open(stringdirectory, 'r').read()
# print(base_path)
# stringtoinsert = open(
#     '/home/atul/Documents/SOFTWARE/OwnPrograms/APPLICATIONS/loginsystem/login_v1/templates/toinsert.html', 'r').read()



@auth.before_app_request
def before_request():
    if current_user.is_authenticated:
        current_user.ping()
        if not current_user.confirmed \
                and request.endpoint \
                and request.blueprint != 'auth' \
                and request.endpoint != 'static':
            return redirect(url_for('auth.unconfirmed'))


def nocache(view):
    @wraps(view)
    def no_cache(*args, **kwargs):
        response = make_response(view(*args, **kwargs))
        response.headers['Last-Modified'] = datetime.now()
        response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, post-check=0, pre-check=0, max-age=0'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '-1'
        return response
    return update_wrapper(no_cache, view)
@auth.route('/unconfirmed')
def unconfirmed():
    if current_user.is_anonymous or current_user.confirmed:
        return redirect(url_for('main.index'))
    return render_template('auth/unconfirmed.html')


@auth.route('/login', methods=['GET', 'POST'])
def login():
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(email=form.email.data).first()
        if user is not None and user.verify_password(form.password.data):
            login_user(user, form.remember_me.data)
            next = request.args.get('next')
            if next is None or not next.startswith('/'):
                next = url_for('main.index')
                session['logged_in'] = True
                session['email']=form.email.data
                session['authenticated'] = False
                if user.confirmed==1:
                    #print(user.is_authenticated)
                    session['authenticated']=True
                # next= url_for('auth.bookslist')
            return redirect(next)
        flash('Invalid username or password.')
    return render_template('auth/login.html', form=form)


@auth.route('/logout')
@login_required
def logout():
    logout_user()
    flash('You have been logged out.')
    session['logged_in'] = False
    return redirect(url_for('main.index'))


@auth.route('/register', methods=['GET', 'POST'])
def register():
    form = RegistrationForm()
    if form.validate_on_submit():
        user = User(email=form.email.data,
                    username=form.username.data,
                    password=form.password.data)
        db.session.add(user)
        db.session.commit()
        token = user.generate_confirmation_token()
        send_email(user.email, 'Confirm Your Account',
                   'auth/email/confirm', user=user, token=token)
        flash('A confirmation email has been sent to you by email.')
        return redirect(url_for('auth.login'))
    return render_template('auth/register.html', form=form)


@auth.route('/confirm/<token>')
@login_required
def confirm(token):
    if current_user.confirmed:
        return redirect(url_for('main.index'))
    if current_user.confirm(token):
        db.session.commit()
        flash('You have confirmed your account. Thanks!')
    else:
        flash('The confirmation link is invalid or has expired.')
    return redirect(url_for('main.index'))


@auth.route('/confirm')
@login_required
def resend_confirmation():
    token = current_user.generate_confirmation_token()
    send_email(current_user.email, 'Confirm Your Account',
               'auth/email/confirm', user=current_user, token=token)
    flash('A new confirmation email has been sent to you by email.')
    return redirect(url_for('main.index'))


@auth.route('/change-password', methods=['GET', 'POST'])
@login_required
def change_password():
    form = ChangePasswordForm()
    if form.validate_on_submit():
        if current_user.verify_password(form.old_password.data):
            current_user.password = form.password.data
            db.session.add(current_user)
            db.session.commit()
            flash('Your password has been updated.')
            return redirect(url_for('main.index'))
        else:
            flash('Invalid password.')
    return render_template("auth/change_password.html", form=form)


@auth.route('/reset', methods=['GET', 'POST'])
def password_reset_request():
    if not current_user.is_anonymous:
        return redirect(url_for('main.index'))
    form = PasswordResetRequestForm()
    if form.validate_on_submit():
        user = User.query.filter_by(email=form.email.data).first()
        if user:
            token = user.generate_reset_token()
            send_email(user.email, 'Reset Your Password',
                       'auth/email/reset_password',
                       user=user, token=token,
                       next=request.args.get('next'))
        flash('An email with instructions to reset your password has been '
              'sent to you.')
        return redirect(url_for('auth.login'))
    return render_template('auth/reset_password.html', form=form)




@auth.route('/reset/<token>', methods=['GET', 'POST'])
def password_reset(token):
    if not current_user.is_anonymous:
        return redirect(url_for('main.index'))
    form = PasswordResetForm()
    if form.validate_on_submit():
        if User.reset_password(token, form.password.data):
            db.session.commit()
            flash('Your password has been updated.')
            return redirect(url_for('auth.login'))
        else:
            return redirect(url_for('main.index'))
    return render_template('auth/reset_password.html', form=form)


@auth.route('/change_email', methods=['GET', 'POST'])
@login_required
def change_email_request():
    form = ChangeEmailForm()
    if form.validate_on_submit():
        if current_user.verify_password(form.password.data):
            new_email = form.email.data
            token = current_user.generate_email_change_token(new_email)
            send_email(new_email, 'Confirm your email address',
                       'auth/email/change_email',
                       user=current_user, token=token)
            flash('An email with instructions to confirm your new email '
                  'address has been sent to you.')
            return redirect(url_for('main.index'))
        else:
            flash('Invalid email or password.')
    return render_template("auth/change_email.html", form=form)


@auth.route('/change_email/<token>')
@login_required
def change_email(token):
    if current_user.change_email(token):
        db.session.commit()
        flash('Your email address has been updated.')
    else:
        flash('Invalid request.')
    return redirect(url_for('main.index'))


@auth.route('/bookslist', methods=['GET', 'POST'])
@login_required
def bookslist():
    user = User.query.filter_by(email=session['email']).first()
    if user.confirmed == 1:
        session['authenticated'] = True
    if session['authenticated']==True:
        try:
            thisuseremail = session['email']
            dataallcol = Book.query.filter_by(email=thisuseremail).all()
            return  render_template('auth/bookslist.html', data=dataallcol)
        except Exception as e:
            return render_template("500.html", error=e)
    return redirect(url_for('main.index'))
    # , TOPIC_DICT=TOPIC_DICT)
    # return render_template("auth/bookslist.html")

@auth.route('/showthebook/<yourbook>', methods=['GET', 'POST'])
@login_required
def showthebook(yourbook):
    return render_template("auth/showthebook.html",yourbook=yourbook)

@auth.route('/profilepage', methods=['GET', 'POST'])
@login_required
def profilepage():
    currentusername=session['email']
    return redirect(url_for('main.user',username=currentusername)) #("user.html",)

@auth.route('/uploadfileprocess', methods=['GET', 'POST'])
@login_required
def uploadfileprocess():
    # print(THIS_DIR)
    if request.method == 'POST':
            # check if the post request has the file part
            if 'file' not in request.files:
                flash('No file part')
                return redirect(url_for('auth.uploadbooks'))
            file = request.files['file']
            # if user does not select file, browser also
            # submit a empty part without filename
            if file.filename == '':
                flash('No selected file')
                return redirect(url_for('auth.uploadbooks'))
            try:

                thisuseremail = session['email']
                thisuseremail = thisuseremail.replace("@", "at")
                thisuseremail = thisuseremail.replace(".", "dot")
                # thisuserdirectory='uploadedimages/' + thisuseremail
                thisuserdirectory = 'static'
                # '''uploadedimages/'
                target1 = os.path.join(THIS_DIR, thisuserdirectory)
                target = os.path.join(target1, thisuseremail)

                # print(target)
                if not os.path.isdir(target):
                    os.mkdir(target)
                # c, conn = connection()
                for file in request.files.getlist("file"):
                    print(file)
                    # flash(thisfilecontent[-10:])


                    filename = file.filename
                    filename = filename.replace(" ", "")
                    thisfileuid = str(uuid.uuid4())[:8]
                    filename_noext=filename[:-4] + thisfileuid
                    session['currentfoldername']=filename_noext
                    book_folder = os.path.join(target, filename_noext)
                    if not os.path.isdir(book_folder):
                        os.mkdir(book_folder)
                    filename = filename[:-4] + thisfileuid + '.pdf'
                    filenamejson = filename[:-4]   + '.json'

                    destinationjson = "/".join([book_folder, filenamejson])
                    filenamecsv = filename[:-4] + '.csv'
                    destinationcsv = "/".join([book_folder, filenamecsv])
                    d = []
                    #
                    d.append({'QUESTION': 'QUESTION', 'ANSWER': 'ANSWER',
                              'TOPIC': 'TOPIC'})

                    d1 = pd.DataFrame(d)

                    d1.to_csv(destinationcsv, encoding='utf-8', index=False)
                    df = pd.read_csv(destinationcsv)
                    df.drop(df.index[0], inplace=True)
                    df.to_csv(destinationcsv, encoding='utf-8', index=False)
                    with open(destinationjson, 'w') as outfile:
                        json.dump([], outfile)

                    print(destinationcsv)
                        # outfile.close()
                    filenamehtml = filename[:-3] + 'html'

                    destination = "/".join([book_folder, filename])
                    destinationhtml = "/".join([book_folder, filenamehtml])
                    # print(destination)

                    # filenamejsaon = 'datap.json'

                    bookname = filename_noext+str(0)
                    file.save(destination)
                    pdf_file = open(destination, 'rb')
                    pdf_reader = PdfFileReader(pdf_file)
                    num_pages=pdf_reader.getNumPages()
                    for x in range(0, num_pages):
                        pdf_writer = PdfFileWriter()
                        pdf_writer.addPage(pdf_reader.getPage(x))
                        pagenum_pdf = destination[:-4] + str(x) + '.pdf'
                        pagenum_html = destination[:-4] + str(x) + '.html'
                        page1 = open(pagenum_pdf, 'wb')
                        pdf_writer.write(page1)
                        page1.close()

                    flash("File uploaded successfully!!")
                    pdf_file.close()

                    thisemail = session['email']
                    thisbook = Book(email=thisemail,
                                    bookname=bookname,
                                    booklocation=book_folder)
                    db.session.add(thisbook)
                    db.session.commit()
                    return redirect(url_for('auth.firstpagedisplay', thisbookname=bookname))
                    # return "XX"
            except Exception as e:
                return render_template("500.html", error=e)
    return redirect(url_for('auth.uploadbooks'))

@auth.route('/firstpage<thisbookname>')
@login_required
def firstpagedisplay(thisbookname):
    try:

        thisuseremail = session['email']
        thisuseremail = thisuseremail.replace("@", "at")
        thisuseremail = thisuseremail.replace(".", "dot")
        thisuserdirectory = 'static'
        target1 = os.path.join(THIS_DIR, thisuserdirectory)
        target = os.path.join(target1, thisuseremail)
        filename_noext = thisbookname[:-1]
        session['currentfoldername']=filename_noext
        book_folder = os.path.join(target, filename_noext)
        filename = thisbookname + '.html'
        totalpage = len(glob.glob1(book_folder, "*.pdf")) -1
        print(totalpage)
        session['book_total_page'] = totalpage
        filenamepdf = thisbookname + '.pdf'
        filenamehtml = thisbookname + '.html'
        session['currentbookname'] = thisbookname
        session['pagenumber'] = 1
        # destination = "/".join([book_folder, filenamepdf])

        pagenum_pdf = "/".join([book_folder, filenamepdf]) #destination[:-4] + str(x) + '.pdf'
        pagenum_html = pagenum_pdf[:-4] + '.html'
        my_file_pdf= Path(pagenum_pdf)
        my_file_html=Path(pagenum_html)
        if my_file_pdf.exists():
            if not my_file_html.exists():
                subprocess.call(["pdf2htmlEX", "--zoom", "3","--dest-dir",book_folder, pagenum_pdf], shell=False)
            # temphtml = filenamehtml #filename_noext + str(x) + '.html'
            # subprocess.call(["mv", temphtml, pagenum_html], shell=False)
                thisfileread = open(pagenum_html, 'r')

                thisfilecontent = thisfileread.read()
                thisfile = open(pagenum_html, 'w')
                thisfilecontentend = thisfilecontent[-17:]
                thisfilecontent = thisfilecontent[:-17]
                # thisfilecontent = u'-'.join(thisfilecontent, stringtoinsert, thisfilecontentend)
                thisfilecontent = thisfilecontent + stringtoinsert + thisfilecontentend
                thisfile.write(smart_str(thisfilecontent))  # .decode('utf-8').encode('utf-8'))
                thisfile.close()

                # print(filename)
                # print(THIS_DIR)


                # book_folder = os.path.join(filename_noext, filename_noext)
                book_folder1 = os.path.join(thisuseremail, filename_noext)
                filename = os.path.join(book_folder1, filename)
                # flash(filename)
                # flash(thisfilecontent)
                # return  thisfilecontent
                # stringtoinsert
                # thisfilecontent
                thispagenumber = 1 #session['pagenumber']
                # totalpage = session['book_total_page']
                pagecount = str(thispagenumber) + '/' + str(totalpage)
                return render_template("auth/showthebook.html", filename=filename,thispagenumber=pagecount)
            else:
                book_folder1 = os.path.join(thisuseremail, filename_noext)
                filename = os.path.join(book_folder1, filename)
                totalpage = len(glob.glob1(book_folder, "*.pdf")) - 1
                print(totalpage)
                session['book_total_page'] = totalpage
                thispagenumber = 1
                pagecount = str(thispagenumber) + '/' + str(totalpage)
                # flash(filename)
                # flash(thisfilecontent)
                # return  thisfilecontent
                # stringtoinsert
                # thisfilecontent
                return render_template("auth/showthebook.html", filename=filename,thispagenumber=pagecount)
        else:
            session['currentbookname'] = filename_noext + str(0)
            thisbookname=session['currentbookname']
            session['pagenumber']=1
            session['book_total_page'] = len(glob.glob1(book_folder, "*.pdf")) - 1
            totalpage=session['book_total_page']
            return redirect(url_for('auth.bookdisplay', thisbookname=thisbookname))
    except Exception as e:
        return render_template("500.html", error=e)


@auth.route('/yourbook<thisbookname>')
@login_required
def bookdisplay(thisbookname):
    try:

        thisuseremail = session['email']
        thisuseremail = thisuseremail.replace("@", "at")
        thisuseremail = thisuseremail.replace(".", "dot")
        thisuserdirectory = 'static'
        target1 = os.path.join(THIS_DIR, thisuserdirectory)
        target = os.path.join(target1, thisuseremail)

        current_foldername = session['currentfoldername']
        # print(current_foldername)
        foldernamelength = len(current_foldername)

        filename_noext = current_foldername #thisbookname[:-1]
        book_folder = os.path.join(target, filename_noext)
        filename = thisbookname + '.html'

        filenamepdf = thisbookname + '.pdf'
        filenamehtml = thisbookname + '.html'
        session['currentbookname'] = thisbookname

        # destination = "/".join([book_folder, filenamepdf])

        pagenum_pdf = "/".join([book_folder, filenamepdf]) #destination[:-4] + str(x) + '.pdf'
        pagenum_html = pagenum_pdf[:-4] + '.html'
        my_file_pdf= Path(pagenum_pdf)
        my_file_html=Path(pagenum_html)
        if my_file_pdf.exists():
            if not my_file_html.exists():
                subprocess.call(["pdf2htmlEX", "--zoom", "3","--dest-dir",book_folder, pagenum_pdf], shell=False)
            # temphtml = filenamehtml #filename_noext + str(x) + '.html'
            # subprocess.call(["mv", temphtml, pagenum_html], shell=False)
                thisfileread = open(pagenum_html, 'r')

                thisfilecontent = thisfileread.read()
                thisfile = open(pagenum_html, 'w')
                thisfilecontentend = thisfilecontent[-17:]
                thisfilecontent = thisfilecontent[:-17]
                # thisfilecontent = u'-'.join(thisfilecontent, stringtoinsert, thisfilecontentend)
                thisfilecontent = thisfilecontent + stringtoinsert + thisfilecontentend
                thisfile.write(smart_str(thisfilecontent))  # .decode('utf-8').encode('utf-8'))
                thisfile.close()

                # print(filename)
                # print(THIS_DIR)


                # book_folder = os.path.join(filename_noext, filename_noext)
                book_folder = os.path.join(thisuseremail, filename_noext)
                filename = os.path.join(book_folder, filename)
                # flash(filename)
                # flash(thisfilecontent)
                # return  thisfilecontent
                # stringtoinsert
                # thisfilecontent
                thispagenumber = session['pagenumber']
                totalpage=session['book_total_page']
                pagecount=str(thispagenumber) +'/'+ str(totalpage)

                return render_template("auth/showthebook.html", filename=filename,thispagenumber = pagecount)
            else:
                book_folder = os.path.join(thisuseremail, filename_noext)
                filename = os.path.join(book_folder, filename)
                # flash(filename)
                # flash(thisfilecontent)
                # return  thisfilecontent
                # stringtoinsert
                # thisfilecontent
                thispagenumber = session['pagenumber']
                totalpage = session['book_total_page']
                pagecount = str(thispagenumber) + '/' + str(totalpage)
                return render_template("auth/showthebook.html", filename=filename,thispagenumber = pagecount)
        else:
            session['currentbookname'] = filename_noext + str(0)
            thisbookname=session['currentbookname']
            session['pagenumber'] =  1
            return redirect(url_for('auth.bookdisplay', thisbookname=thisbookname))
    except Exception as e:
        return render_template("500.html", error=e)
@auth.route('/makeflashcard', methods=['GET', 'POST'])
@login_required
def makeflashcard():

    request_json = request.get_json()
    data = request_json
    #
    #
    # d = []
    # for x in range(0,len(data)):
    #     d.append({'QUESTION':data[x]['question'],'ANSWER':data[x]['answer'],'TOPIC':data[x]['topic']})
    #
    # d1=pd.DataFrame(d)
    # d1.to_csv("master_panda.csv", encoding='utf-8')
    # csvfile.close()
    # thisfiletemp = open("temp1.txt", 'w')
    # for in range(0,len(data)):
    # thisfiletemp.write(smart_str(data[0]['question']))
    # flash("XX")
    # thisfiletemp.close()


    thisuseremail = session['email']
    thisuseremail = thisuseremail.replace("@", "at")
    thisuseremail = thisuseremail.replace(".", "dot")
    thisuserdirectory = 'static'
    target1 = os.path.join(THIS_DIR, thisuserdirectory)
    target = os.path.join(target1, thisuseremail)
    current_foldername = session['currentfoldername']
    filename_noext = current_foldername  # thisbookname[:-1]
    book_folder = os.path.join(target, filename_noext)
    # filename = session['currentbookname']
    filename=session['currentfoldername']
    filename = filename + '.json'
    # filenamejsaon = 'datap.json'
    destinationjson = "/".join([book_folder, filename])
    filenamecsv = session['currentfoldername'] + '.csv'
    destinationcsv = "/".join([book_folder, filenamecsv])
    d= pd.read_csv(destinationcsv)
    # for x in range(0,len(data)):
        # d.append({'QUESTION':data[x]['question'],'ANSWER':data[x]['answer'],'TOPIC':data[x]['topic']})

    d1 = pd.DataFrame(d)
    d1=d1.drop_duplicates()
    d1.to_csv(destinationcsv, mode='w', encoding='utf-8', index=False)

    # with open(destinationjson, mode='r', encoding='utf-8') as feedsjson:
    with open(destinationjson, 'r') as justread:
        previousqa = json.load(justread)
    with open(destinationjson, 'w') as outfile:
        previousqa.append(data)
        json.dump(previousqa, outfile)
    # subprocess.call(["mv", filename, destinationjson], shell=False)
    # thisfile = open(destinationhtml, 'w')
    # # thisfilecontent = thisfile.read()
    # thisfile.write(name.decode('utf-8').encode('utf-8'))
    # thisfile.write(smart_str(name))
    # flash(name)
    # thisfile.close()
    return redirect(url_for('auth.makeppt'))

        # render_template("500.html", error=e) <int:id
    #
@auth.route('/showflashcardspq', methods=['GET', 'POST'])
@login_required
def showflashcardspq():
    data = session['flashcards']
    len(data)
    text = request.form['thisqnumber']
    urlforppt = session['urlforppt']
    print(text)
    if(len(text)>0):
        qnumbershow=int(text)
        qnumber=qnumbershow-1
        print(qnumber)

        if qnumber >= 0 and qnumber < len(data):
            datatosend = data[qnumber]
            session['qnumber']=qnumber
            return render_template('flashcard.html', datatosend=datatosend, urlforppt=urlforppt, qnumber=qnumber)
    qnumber=int(0) #session['qnumber']
    datatosend = data[qnumber]
    return render_template('flashcard.html', datatosend=datatosend, urlforppt=urlforppt, qnumber=qnumber)

@auth.route('/showflashcard2<int:qnumber>', methods=['GET', 'POST'])
@login_required
def showflashcard2(qnumber):
    data = session['flashcards']

    qnumber = qnumber
    urlforppt = session['urlforppt']
    if request.method == "POST":
        if request.form['action'] == 'previousqa':
            qnumber = qnumber - 1
            if qnumber >= 0:

                datatosend= data[qnumber]
                session['qnumber'] = qnumber
                return render_template('flashcard.html',datatosend=datatosend ,urlforppt=urlforppt,qnumber=qnumber )

            else:
                qnumber = qnumber + 1
                datatosend = data[qnumber]
                session['qnumber'] = qnumber
                return render_template('flashcard.html', datatosend=datatosend, urlforppt=urlforppt, qnumber=qnumber)


        elif request.form['action'] == 'nextqa':
            qnumber = qnumber + 1
            if qnumber < len(data):
                datatosend = data[qnumber]
                session['qnumber'] = qnumber
                return render_template('flashcard.html', datatosend=datatosend, urlforppt=urlforppt, qnumber=qnumber)
            else:
                qnumber = qnumber - 1
                datatosend = data[qnumber]
                session['qnumber'] = qnumber
                return render_template('flashcard.html', datatosend=datatosend, urlforppt=urlforppt,
                                       qnumber=qnumber)
    datatosend = data[qnumber]
    return render_template('flashcard.html', qnumber=qnumber,datatosend=datatosend, urlforppt=urlforppt)

@auth.route('/gotopage', methods=['GET', 'POST'])
@login_required
def gotopage():

    current_foldername = session['currentfoldername']



    text = request.form['thispagenumber']
    print(text)
    if (len(text) > 0):
        pnumbertoshow = int(text)
        pnumber = pnumbertoshow - 1

        if pnumber >= 0 and pnumber < session['book_total_page']:
            thispagenumberpdf = current_foldername + str(pnumber)
            session['pagenumber'] = pnumber + 1
            return redirect(url_for('auth.bookdisplay', thisbookname=thispagenumberpdf))
    thispagenumberpdf = current_foldername + str(session['pagenumber'])
    return redirect(url_for('auth.bookdisplay', thisbookname=thispagenumberpdf))


@auth.route('/shownextpage', methods=['GET', 'POST'])
@login_required
def shownextpage():
    current_foldername=session['currentfoldername']
    # print(current_foldername)
    foldernamelength=len(current_foldername)
    thisbookname = session['currentbookname']
    lastpagenumber=thisbookname[foldernamelength:]
    thispagenumber=int(lastpagenumber) + 1
    session['pagenumber'] = thispagenumber + 1
    thispagenumberpdf = current_foldername + str(thispagenumber)
    # session['pagenumber'] = thispagenumber + 1
    # print(str(session['pagenumber']))
    # print

    return redirect(url_for('auth.bookdisplay', thisbookname=thispagenumberpdf))

@auth.route('/showcurrentpage', methods=['GET', 'POST'])
@login_required
def showcurrentpage():
    thisbookname=session['currentbookname']
    # session['pagenumber'] = thispagenumber + 1
    # lastpagenumber = thisbookname[-1:]
    # thispagenumber = ord(lastpagenumber) - 1
    # thispagenumberpdf = thisbookname[:-1] + str(thispagenumber)
    # thispagenumber=str(session['pagenumber'])
    # print(str(session['pagenumber']))

    # print(thispagenumber)
    return redirect(url_for('auth.bookdisplay', thisbookname=thisbookname))

@auth.route('/showpreviouspage', methods=['GET', 'POST'])
@login_required
def showpreviouspage():
    current_foldername=session['currentfoldername']
    foldernamelength=len(current_foldername)
    thisbookname = session['currentbookname']
    lastpagenumber=thisbookname[foldernamelength:]
    thispagenumber=int(lastpagenumber) - 1
    session['pagenumber'] = thispagenumber+1
    thispagenumberpdf = current_foldername + str(thispagenumber)

    # print(str(session['pagenumber']))
    # print(thispagenumberpdf)
    return redirect(url_for('auth.bookdisplay', thisbookname=thispagenumberpdf))


@auth.route('/showflashcard1<int:qnumber>', methods=['GET', 'POST'])
@login_required
def showflashcard1(qnumber):
    data = session['flashcards']
    question=data[qnumber]['question']
    answer = data[qnumber]['answer']
    qnumber = qnumber
    urlforppt = session['urlforppt']
    if request.method == "POST":
        if request.form['action'] == 'previousqa':
            qnumber = qnumber - 1
            if qnumber >= 0:
                question = data[qnumber]['question']
                answer = data[qnumber]['answer']
                return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                       question=question)
            else:
                qnumber = qnumber + 1
                question = data[qnumber]['question']  + " !!!THIS IS THE FIRST QUESTION!!!"
                answer = data[qnumber]['answer']  + " !!!THIS IS THE FIRST ANSWER!!!"
                return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                       question=question)
        elif request.form['action'] == 'nextqa':
            qnumber = qnumber + 1
            if qnumber < len(data):
                question = data[qnumber]['question']
                answer = data[qnumber]['answer']
                return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                       question=question)
            else:
                qnumber = qnumber - 1
                question = data[qnumber]['question'] + " !!!THIS IS THE LAST QUESTION!!!"
                answer = data[qnumber]['answer'] + " !!!THIS IS THE LAST ANSWER!!!"
                return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                       question=question)
    return render_template('flashcard.html', qnumber=qnumber,urlforppt = urlforppt, answer=answer,question=question)

@auth.route('/showflashcardnext<int:qnumber>', methods=['GET', 'POST'])
@login_required
def showflashcardnext(qnumber):
    qnumber=qnumber+1
    urlforppt = session['urlforppt']
    data = session['flashcards']
    if qnumber < len(data):
        question=data[qnumber]['question']
        answer = data[qnumber]['answer']
        qnumber = qnumber
        return render_template('flashcard.html', qnumber=qnumber,urlforppt = urlforppt, answer=answer,question=question)
    qnumber=qnumber-1
    question = data[qnumber]['question']
    answer = data[qnumber]['answer']
    return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer, question=question)
@auth.route('/showflashcardprev<int:qnumber>', methods=['GET', 'POST'])
@login_required
def showflashcardprev(qnumber):
    data = session['flashcards']
    urlforppt = session['urlforppt']
    if request.form['action'] == 'previousqa':
        qnumber = qnumber-1
        if qnumber >= 0:
            question = data[qnumber]['question']
            answer = data[qnumber]['answer']
            return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                   question=question)
        else:
            qnumber = qnumber+1
            question = data[qnumber]['question']
            answer = data[qnumber]['answer']
            return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                   question=question)


    elif request.form['action'] == 'nextqa':
        qnumber=qnumber+1
        if qnumber < len(data):
            question = data[qnumber]['question']
            answer = data[qnumber]['answer']
            return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer, question=question)
        else:
            qnumber=qnumber-1
            question = data[qnumber]['question']
            answer = data[qnumber]['answer']
            return render_template('flashcard.html', qnumber=qnumber, urlforppt=urlforppt, answer=answer,
                                   question=question)




@auth.route('/showflashcard', methods=['GET', 'POST'])
@login_required
def showflashcard():

    thisuseremail = session['email']
    thisuseremail = thisuseremail.replace("@", "at")
    thisuseremail = thisuseremail.replace(".", "dot")
    thisuserdirectory = 'static'
    target1 = os.path.join(THIS_DIR, thisuserdirectory)
    target = os.path.join(target1, thisuseremail)
    filename = session['currentbookname']
    filenamepptx = filename[:-4] + 'pptx'
    urlforppt = os.path.join(thisuseremail, filenamepptx)
    filename = filename[:-4] + 'json'
    # filename="spring.json"
    destinationjson = "/".join([target, filename])
    destinationpptx = "/".join([target, filenamepptx])
    json_url = destinationjson
    # os.path.join(SITE_ROOT, destinationjson)
    data = json.load(open(json_url))

    # data=json.JSONEncoder(data)
    listofflashcard=[]

    data_length = len(data)
    for datap in range(0, data_length):
        flashcarddata = Flashcardqa()
        flashcarddata.questionnum = smart_str(data[datap]['qanumber'])
        flashcarddata.answer = smart_str(data[datap]['answer'])
        flashcarddata.question = smart_str(data[datap]['question'])
        listofflashcard.append(flashcarddata)
        # print(flashcarddata.question)
    urlforfc="/"
    # urlforppt="/"
    data=session['flashcards']
    data = json.dumps(data)
    return render_template('flashcard.html',listofflashcard=listofflashcard, data=data,urlforfc=urlforfc,urlforppt=urlforppt)



@auth.route('/makeppt', methods=['GET', 'POST'])
@login_required
def makeppt():
    ###########
    thisuseremail = session['email']
    # thisuseremail = thisuseremail.replace("@", "at")
    # thisuseremail = thisuseremail.replace(".", "dot")
    # thisuserdirectory = 'static'
    # target1 = os.path.join(THIS_DIR, thisuserdirectory)
    # target = os.path.join(target1, thisuseremail)

    current_foldername = session['currentfoldername']
    thisbookname = session['currentbookname']
    # print(current_foldername)
    # foldernamelength = len(current_foldername)
    # thisbookname = session['currentbookname']
    # lastpagenumber = thisbookname[foldernamelength:]
    # thispagenumber = int(lastpagenumber) + 1
    # thispagenumberpdf = current_foldername + str(thispagenumber)

    filename_noext = current_foldername  # thisbookname[:-1]

    # filename = thisbookname + '.html'

    filenamepdf = thisbookname + '.pdf'
    # filenamehtml = thisbookname + '.html'


    # destination = "/".join([book_folder, filenamepdf])

     # destination[:-4] + str(x) + '.pdf'

    ##########
    thisuseremail = session['email']
    thisuseremail = thisuseremail.replace("@", "at")
    thisuseremail = thisuseremail.replace(".", "dot")
    thisuserdirectory = 'static'
    target1 = os.path.join(THIS_DIR, thisuserdirectory)
    target = os.path.join(target1, thisuseremail)
    book_folder = os.path.join(target, filename_noext)
    # filename = session['currentbookname']
    filename=session['currentfoldername']
    # print(filename)
    # pagenum_pdf = "/".join([book_folder, filenamepdf])
    # pagenum_html = pagenum_pdf[:-4] + '.html'
    filenamepptx=filename + '.pptx'
    # filename = filename + '.json'
    filenamecsv = session['currentfoldername'] + '.csv'
    # filenamejsaon = 'datap.json'
    # destinationjson = "/".join([book_folder, filename])
    destinationpptx = "/".join([book_folder, filenamepptx])
    destinationcsv = "/".join([book_folder, filenamecsv])
    # print(destinationjson)
    data_panda = pd.read_csv(destinationcsv)
    data=[]
    for xx in range(len(data_panda)):
        data.append({'question': data_panda['QUESTION'][xx],'answer':data_panda['ANSWER'][xx], 'topic':data_panda['TOPIC'][xx]})
    # session['flashcards_panda']=data_panda
    # data = json.load(open(destinationjson))
    # data=data_panda.to_json(orient='records')

    # data=data_panda
    # data.append('AA')
    session['flashcards'] = data
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]
    listofflashcard=[]
    data_length = len(data)
    for datap in range(0, data_length):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = smart_str(data[datap]['topic'])
        tf = body_shape.text_frame
        tf.text = smart_str(data[datap]['question'])
        p = tf.add_paragraph()
        p.text = smart_str(data[datap]['answer'])
        p.level = 1
        # print(listofflashcard[datap].answer)
    urlforppt1 = os.path.join(thisuseremail, filename_noext)
    urlforppt = os.path.join(urlforppt1, filenamepptx)
    prs.save(destinationpptx)
    session['urlforppt'] = urlforppt
    # session['flashcardqalist']=listofflashcard
    return redirect(url_for('auth.showflashcard2',qnumber=0))



@auth.route('/process', methods=['GET', 'POST'])
@login_required
def process():
    # try:
    request_json = request.get_json()
    data = request_json['email']

    name = request_json['name']

    # return email
    name = '<!DOCTYPE html><html xmlns="http://www.w3.org/1999/xhtml">' + name + '</html>'
    # name =name.encode('utf-8')
    newName = name
#################
    thisuseremail = session['email']
    thisuseremail = thisuseremail.replace("@", "at")
    thisuseremail = thisuseremail.replace(".", "dot")
    thisuserdirectory = 'static'
    target1 = os.path.join(THIS_DIR, thisuserdirectory)
    target = os.path.join(target1, thisuseremail)
    thisbookname=session['currentbookname']
    filename_noext = session['currentfoldername']
    book_folder = os.path.join(target, filename_noext)
    filename = thisbookname + '.html'

    filenamepdf = thisbookname + '.pdf'
    filenamehtml = thisbookname + '.html'
    # destination = "/".join([book_folder, filenamepdf])

    pagenum_pdf = "/".join([book_folder, filenamepdf])  # destination[:-4] + str(x) + '.pdf'
    pagenum_html = pagenum_pdf[:-4] + '.html'

    filenamecsv = session['currentfoldername'] + '.csv'
    destinationcsv = "/".join([book_folder, filenamecsv])
    d = []  # pd.read_csv(destinationcsv)
    for x in range(0, len(data)):
        d.append({'QUESTION': data[x]['question'], 'ANSWER': data[x]['answer'], 'TOPIC': data[x]['topic']})

    d1 = pd.DataFrame(d)

    d1.to_csv(destinationcsv, mode='a', header=False, encoding='utf-8', index=False)
    ###################################
    # print(thisbookname)
    # print(pagenum_html)

    # thisuseremail = session['email']
    # thisuseremail = thisuseremail.replace("@", "at")
    # thisuseremail = thisuseremail.replace(".", "dot")
    # thisuserdirectory = 'static'
    # target1 = os.path.join(THIS_DIR, thisuserdirectory)
    # target = os.path.join(target1, thisuseremail)
    # filename = session['currentbookname']

    # destinationhtml =pagenum_html# "/".join([target, filename])
    jsonify({'name': name})

    thisfile = open(pagenum_html, 'w')
    # # thisfilecontent = thisfile.read()
    # thisfile.write(name.decode('utf-8').encode('utf-8'))
    thisfile.write(smart_str(name))
    # flash(name)
    # print(name)
    thisfile.close()
    # thisbookname=session['currentbookname']
    # dataallcol = Book.query.filter_by(bookname=thisbookname).first()
    # request.form['name']
    # subprocess.call(["echo", name[::-1], ">>","newhtml.html"], shell=False)
    # if name and email:
    #     return "XX"
        # jsonify({'name' : name})
    # print(name)
    return redirect(url_for('auth.showcurrentpage'))
    #     # return "XX"
    # except Exception as e:
    #     return render_template("500.html", error=e)


@auth.route('/uploadwebaddress', methods=['GET', 'POST'])
@login_required
def uploadwebaddress():
    if request.method == 'POST':
        # check if the post request has the file part
        weblinkname = request.form['web_add_name']
        linkpdfname = request.form['web_name']
        linkpdfname = linkpdfname.replace(" ", "")
        # val = URLValidator(verify_exists=True)

        if not validators.url(weblinkname):
        # except ValidationError, e:
            flash('Invalid URL')
            return redirect(url_for('auth.uploadbooks'))
        # weblinkname=weblinkname.linkpdfname.replace(" ", "")
        if len(linkpdfname)<1:
            flash('Invalid file name')
            return redirect(url_for('auth.uploadbooks'))

        try:
            if request.form['web_add_name'] and request.form['web_name']:
                weblinkname = request.form['web_add_name']
                linkpdfname = request.form['web_name']
                linkpdfname=linkpdfname.replace(" ", "")
                thisfileuid=str(uuid.uuid4())[:8]
                filename_noext = linkpdfname + thisfileuid
                linkpdfnamepdf = linkpdfname + thisfileuid +'.pdf'

                thisuseremail = session['email']
                thisuseremail = thisuseremail.replace("@", "at")
                thisuseremail = thisuseremail.replace(".", "dot")
                # thisuserdirectory='uploadedimages/' + thisuseremail
                thisuserdirectory = 'static'
                # '''uploadedimages/'
                target1 = os.path.join(THIS_DIR, thisuserdirectory)
                target = os.path.join(target1, thisuseremail)
                if not os.path.isdir(target):
                    os.mkdir(target)
                session['currentfoldername'] = filename_noext
                print(filename_noext)
                book_folder = os.path.join(target, filename_noext)
                if not os.path.isdir(book_folder):
                    os.mkdir(book_folder)
                # c, conn = connection()
                filename = linkpdfnamepdf
                filenamehtml = filename[:-3] + 'html'
                filenamejson = filename[:-3]+ 'json'
                filenamecsv = filename[:-3]+ 'csv'
                destinationcsv = "/".join([book_folder, filenamecsv])
                d = []
                #
                d.append({'QUESTION': 'QUESTION', 'ANSWER': 'ANSWER',
                          'TOPIC': 'TOPIC'})

                d1 = pd.DataFrame(d)
                d1.drop(d1.index[0])
                d1.to_csv(destinationcsv, encoding='utf-8', index=False)

                destinationjson = "/".join([book_folder, filenamejson])
                with open(destinationjson, 'w') as outfile:
                    json.dump([], outfile)
                destination = "/".join([book_folder, filename])
                destinationhtml = "/".join([book_folder, filenamehtml])
                bookname = filename_noext + str(0)
                # destination = "/".join([target, filename])
                # destinationhtml = "/".join([target, filenamehtml])
                # # pdfkit.from_url(weblinkname, destination)
                subprocess.call(["wget", "-k", weblinkname,"-O", destinationhtml], shell=False)
                soup = BeautifulSoup(open(destinationhtml), "html.parser")
                new_link = soup.new_tag("link")
                new_link.attrs["href"] = "removeheader.css"
                new_link.attrs["rel"] = "stylesheet"
                soup.head.insert(0, new_link)
                for m in soup.find_all('a'):
                    m.unwrap()
                    # replaceWithChildren()
                # soup = BeautifulSoup(page)
                with open(destinationhtml, "w") as file:
                    file.write(smart_str(soup))
                file.close()
                # subprocess.call(["cp", filenamehtml, "temp.html"], shell=False)
                thiswebpagepdfparam = "--print-to-pdf=" + destination
                #
                subprocess.call(["google-chrome", "--headless", "--disable-gpu", thiswebpagepdfparam, destinationhtml], shell=False)
                pdf_file = open(destination, 'rb')
                pdf_reader = PdfFileReader(pdf_file)
                num_pages = pdf_reader.getNumPages()
                for x in range(0, num_pages):
                    pdf_writer = PdfFileWriter()
                    pdf_writer.addPage(pdf_reader.getPage(x))
                    pagenum_pdf = destination[:-4] + str(x) + '.pdf'
                    pagenum_html = destination[:-4] + str(x) + '.html'
                    page1 = open(pagenum_pdf, 'wb')
                    pdf_writer.write(page1)
                    page1.close()

                flash("File uploaded successfully!!")
                pdf_file.close()
                thisemail = session['email']
                thisbook = Book(email=thisemail,
                                bookname=bookname,
                                booklocation=book_folder)
                db.session.add(thisbook)
                db.session.commit()
                return redirect(url_for('auth.bookdisplay', thisbookname=bookname))
                # subprocess.call(["pdf2htmlEX", "--zoom","3",  filename], shell=False)
                # # page = urllib2.urlopen(filenamehtml).read()
                #
                # subprocess.call(["mv", filenamehtml, destinationhtml], shell=False)
                # subprocess.call(["rm", filename], shell=False)
                # # subprocess.call(["mv", filenamehtml, destinationhtml], shell=False)
                # thisfileread = open(destinationhtml, 'r')
                #
                # thisfilecontent = thisfileread.read()
                # thisfile = open(destinationhtml, 'w')
                # # re.sub("home*?html", "", thisfile)
                # thisfilecontentend = thisfilecontent[-17:]
                # thisfilecontent = thisfilecontent[:-17]
                # # thisfilecontent = u'-'.join(thisfilecontent, stringtoinsert, thisfilecontentend)
                # thisfilecontent = thisfilecontent + stringtoinsert + thisfilecontentend
                # thisfile.write(smart_str(thisfilecontent))
                # # thisfile.write(thisfilecontent.encode('utf8') + '\n')
                # # thisfile.write(thisfilecontent.decode('utf-8').encode('utf-8'))
                # thisfile.close()
                # bookname =  filename[:-4]
                #
                # thisemail = session['email']
                # thisbook = Book(email=thisemail,
                #                 bookname=bookname,
                #                 booklocation=destinationhtml)
                # db.session.add(thisbook)
                # db.session.commit()
                # flash("File uploaded successfully!!")
                # return redirect(url_for('auth.bookdisplay', thisbookname=bookname))
        except Exception as e:
                return render_template("500.html", error=e)
    return redirect(url_for('auth.uploadbooks'))


@auth.route('/uploadbooks', methods=['GET', 'POST'])
@login_required
def uploadbooks():
    user = User.query.filter_by(email=session['email']).first()
    if user.confirmed == 1:
        session['authenticated'] = True
    if session['authenticated'] == True:
        return render_template("auth/uploadbooks.html")
    return redirect(url_for('main.index'))

@auth.route('/termsnconditions')
def termsnconditions():
    return render_template('termsnconditions.html')

@auth.route('/who')
def who():
    return render_template('who.html')

@auth.route('/why')
def why():
    return render_template('why.html')

@auth.route('/how')
def how():
    return render_template('how.html')

@auth.route('/homepage')
def homepage():
    return render_template('homepage.html')

@auth.route('/delete/<int:id>', methods=['POST'])
def remove(id):
    booktoremove = Book.query.filter_by(id=id).first()
    print(booktoremove.booklocation)
    strintoremove=booktoremove.booklocation #[:-4]+'*'
    subprocess.call(["rm", "-r", strintoremove], shell=False)
    # strintoremovepdf=strintoremove+'pdf'
    # strintoremovejson = strintoremove + 'json'
    # subprocess.call(["rm", "-r", booktoremove.booklocation], shell=False)

    #subprocess.call("rm", "-r", str(booktoremove.booklocation))
    db.session.delete(booktoremove)
    db.session.commit()

    #object = Object.query.get_or_404(id)
    #delete(object)
    return redirect(url_for('auth.bookslist'))

